#!/usr/bin/env python3
import os
import logging
import sys
import random
import shutil

from tkinter import messagebox
from tkinter import *
import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from pptxtopdf import convert


class CertificateCreator:
    def __init__(self, windows):
        self.window = windows
        self.window.title('Создание сертификата')
        self.window.geometry('450x300')

        self.name = tk.StringVar()
        self.price = tk.StringVar()
        self.buyer = tk.StringVar()

        self.setup_ui()
        self.setup_logging()

    @staticmethod
    def setup_logging():
        """Настраивает логирование"""
        logging.basicConfig(
            level=logging.INFO,
            filename='program.log',
            encoding='WINDOWS-1251',
            filemode='a',
            format='%(asctime)s - %(levelname)s - %(message)s',
            datefmt='%d.%m.%Y - %H:%M'
        )
        logger = logging.getLogger()
        stream_handler = logging.StreamHandler(sys.stdout)
        stream_handler.setLevel(logging.INFO)
        stream_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        logger.addHandler(stream_handler)

    def setup_ui(self):
        """Рисует графику с помощью tkinter"""
        frame = Frame(self.window)
        frame.pack()

        ttk.Label(frame, text='ФИО').pack(padx=100, pady=5)
        ttk.Entry(frame, textvariable=self.name, width=30).pack(padx=100, pady=5)

        ttk.Label(frame, text='СУММА').pack(padx=100, pady=5)
        ttk.Entry(frame, textvariable=self.price, width=30).pack(padx=100, pady=5)

        ttk.Label(frame, text='КОМУ ПРОДАНО').pack(padx=100, pady=5)
        ttk.Entry(frame, textvariable=self.buyer, width=30).pack(padx=100, pady=5)

        btn_create = tk.Button(
            self.window, text='Создать сертификат',
            command=self.generate_certificate)
        btn_create.pack(pady=5)

        exit_button = tk.Button(self.window, text="Выход", command=self.window.destroy)
        exit_button.pack(pady=5)

    @staticmethod
    def get_random_number():
        """Получает случайное число до 6 знаков."""
        return random.randint(100000, 999999)

    @staticmethod
    def check_input_file(input_pptx):
        """Проверяет существование файлов в папке с программой."""
        if not os.path.exists(input_pptx):
            messagebox.showinfo(
                message=f'Файл «{input_pptx}» не найден в папке с программой.'
            )
            return False
        return True

    def generate_certificate(self):
        """Генерирует сертификат."""
        input_pptx = './template/Сертификат_шаблон.pptx'
        out_pptx = './pptx/Сертификат.pptx'
        out_pptx_path = './pptx/'
        out_pdf = f'./pdf/'
        pdf_file = f'./pdf/Сертификат.pdf'
        price_value = str(self.price.get())
        name_value = str(self.name.get())
        buyer_value = str(self.buyer.get())
        prs = Presentation(input_pptx)


        for folder in [out_pptx_path, out_pdf]:
            if os.path.exists(folder):
                try:
                    shutil.rmtree(folder)
                except Exception as e:
                    messagebox.showerror(
                        "Ошибка",
                        f"Не удалось удалить папку {folder}: {str(e)}"
                    )
                    return

        prs = Presentation(input_pptx)

        if not self.check_input_file(input_pptx):
            return

        if len(price_value) > 6:
            messagebox.showerror(
                "Ошибка",
                "Цена слишком большая (максимум 100000)"
            )
            return

        if len(price_value) <= 0:
            messagebox.showerror(
                "Ошибка",
                "Цена должна быть больше 0"
            )
            logging.error("Цена должна быть больше 0")
            return

        replacements = {
            'price': str(price_value),
            'name': name_value,
            'serial': str(self.get_random_number()),
        }

        slide = prs.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        run.text = run.text.replace(key, value)

        # Создаем папки заново
        os.makedirs(out_pptx_path, exist_ok=True)
        os.makedirs(out_pdf, exist_ok=True)
        
        prs.save(out_pptx)

        try:
            self.convert_pptx_to_pdf(out_pptx, out_pdf, pdf_file)
        except Exception as e:
            messagebox.showerror(
                'Ошибка',
                f'Не удалось конвертировать: {e}'
            )

        message = (
            f'Сертификат №: {replacements["serial"]}; '
            f'Именной: {replacements["name"]}; '
            f'Кому: {buyer_value}; '
            f'Цена: {replacements["price"]} ₽'
        )
        messagebox.showinfo(message=message)
        logging.info(message)
        self.window.destroy()

    @staticmethod
    def convert_pptx_to_pdf(out_pptx, out_pdf, pdf_file):
        """Конвертирует файл pptx в pdf."""
        try:
            if os.path.exists(pdf_file):
                os.remove(pdf_file)
            convert(out_pptx, out_pdf)

        except Exception as e:
            messagebox.showerror(
                "Ошибка",
                f"Не удалось конвертировать файл: {str(e)}"
            )

    def main(self):
        self.window.mainloop()


if __name__ == '__main__':
    window = Tk()
    app = CertificateCreator(window)
    app.main()
